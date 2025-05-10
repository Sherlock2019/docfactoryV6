import streamlit as st
from docx import Document
from pptx import Presentation
from docx.shared import Inches
import pandas as pd
from io import BytesIO
import os
from datetime import date
from tempfile import NamedTemporaryFile
from pathlib import Path
import json
import re

st.set_page_config(page_title="📄 Rackspace DocFactory", layout="wide")
st.title("📄📊 Rackspace Documentation Factory")

# Setup directories
TEMPLATE_DIR = Path("templates")
TEMPLATE_INDEX = TEMPLATE_DIR / "templates_index.json"
TEMPLATE_DIR.mkdir(exist_ok=True)

# Load or initialize template index
if TEMPLATE_INDEX.exists():
    template_index = json.loads(TEMPLATE_INDEX.read_text())
else:
    template_index = {
        "Solution Proposal": [],
        "Cloud Readiness Assessment": [],
        "Statement of Work": []
    }

today = date.today().strftime("%Y%m%d")

# Sidebar – Template Upload
st.sidebar.markdown("## 🧰 Template Management")
with st.sidebar.expander("➕ Add or Update Template"):
    uploaded = st.file_uploader("Upload .dot / .dotx / .pptx Template", type=["dot", "dotx", "pptx"])
    doc_type = st.selectbox("Assign to Document Type", list(template_index.keys()))
    display_name = st.text_input("Template Display Name")

    if st.button("💾 Save Template") and uploaded and display_name:
        ext = uploaded.name.split(".")[-1]
        filename = f"{doc_type.replace(' ', '_')}_{display_name.replace(' ', '_')}.{ext}"
        path = TEMPLATE_DIR / filename
        with open(path, "wb") as f:
            f.write(uploaded.read())
        entry = {"name": display_name, "file": str(path)}
        template_index[doc_type] = [e for e in template_index[doc_type] if e["name"] != display_name]
        template_index[doc_type].append(entry)
        TEMPLATE_INDEX.write_text(json.dumps(template_index, indent=2))
        st.success(f"✅ Template saved as {filename}")

# Select template to use
doc_type = st.selectbox("📄 Select Document Type", list(template_index.keys()))
customer_name = st.text_input("👤 Customer Name")
template_options = [t["name"] for t in template_index[doc_type]]
template_name = st.selectbox("📑 Select Template", template_options)

# Template logic
template_path = None
for t in template_index[doc_type]:
    if t["name"] == template_name:
        template_path = t["file"]

TEXT_ONLY_PLACEHOLDERS = {"CUSTOMER_NAME", "CITY NAME", "SA-NAME", "SA_EMAIL", "RAX_TEAM", "PARTNER_NAME"}

if template_path and customer_name:
    is_docx = template_path.endswith((".dotx", ".docx", ".dot"))
    is_pptx = template_path.endswith(".pptx")
    uploads = {}

    # Extract placeholders
    text_blocks = []
    if is_docx:
        doc = Document(template_path)
        text_blocks = [p.text for p in doc.paragraphs]
    elif is_pptx:
        prs = Presentation(template_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_blocks.append(shape.text)

    raw_placeholders = re.findall(r"\{[^}]+\}", "\n".join(text_blocks))
    placeholders = list(dict.fromkeys([f"{{{ph.strip('{}').strip()}}}" for ph in raw_placeholders]))

    # Step 1: Manual Text Inputs
    st.markdown("### ✏️ Fill in Required Fields")
    for ph in placeholders:
        base = ph.strip("{}").strip()
        if base in TEXT_ONLY_PLACEHOLDERS:
            val = st.text_input(f"{ph}", key=f"text_{base}")
            if val.strip():
                uploads[ph] = val.strip()

    # Step 2: Upload or Enter for Other Fields
    st.markdown("### 📎 Upload or Enter Remaining Content")
    for ph in placeholders:
        base = ph.strip("{}").strip()
        if base not in TEXT_ONLY_PLACEHOLDERS:
            col1, col2 = st.columns(2)
            with col1:
                file = st.file_uploader(f"📎 Upload for {ph}", type=["txt", "docx", "pptx", "xlsx", "jpg", "png"], key=f"file_{base}")
            with col2:
                val = st.text_area(f"✏️ Or enter value for {ph}", key=f"text_{base}")
            if file:
                ext = file.name.lower().split(".")[-1]
                if ext in ["jpg", "jpeg", "png"]:
                    uploads[ph] = BytesIO(file.read())
                elif ext == "xlsx":
                    df = pd.read_excel(file)
                    uploads[ph] = df
                elif ext == "docx":
                    d = Document(file)
                    uploads[ph] = "\n".join(p.text for p in d.paragraphs)
                elif ext == "pptx":
                    p = Presentation(file)
                    uploads[ph] = "\n".join(shape.text for slide in p.slides for shape in slide.shapes if hasattr(shape, "text"))
                elif ext == "txt":
                    uploads[ph] = file.read().decode("utf-8")
            elif val.strip():
                uploads[ph] = val.strip()

    # Step 3: Generate Output
    missing = [ph for ph in placeholders if ph not in uploads]
    if st.button("🛠️ Generate Document"):
        if missing:
            st.warning(f"⚠️ Missing placeholders: {', '.join(missing)}")
        else:
            final_filename = f"{customer_name}_{doc_type.replace(' ', '_')}_{today}"
            buffer = BytesIO()

            if is_docx:
                doc = Document(template_path)
                for para in doc.paragraphs:
                    for ph, val in uploads.items():
                        if ph in para.text:
                            para.text = para.text.replace(ph, "")
                            if isinstance(val, BytesIO):
                                val.seek(0)
                                with NamedTemporaryFile(delete=False, suffix=".png") as tmp:
                                    tmp.write(val.read())
                                    tmp.flush()
                                    new_para = doc.add_paragraph()
                                    new_para.add_run().add_picture(tmp.name, width=Inches(4))
                                    os.unlink(tmp.name)
                            elif isinstance(val, pd.DataFrame):
                                table = doc.add_table(rows=1, cols=len(val.columns))
                                hdr_cells = table.rows[0].cells
                                for i, col in enumerate(val.columns):
                                    hdr_cells[i].text = col
                                for _, row in val.iterrows():
                                    row_cells = table.add_row().cells
                                    for i, cell in enumerate(row):
                                        row_cells[i].text = str(cell)
                            else:
                                para.add_run(str(val))
                doc.save(buffer)
                st.success("✅ DOCX generated!")
                st.download_button("📥 Download DOCX", buffer.getvalue(), file_name=final_filename + ".docx")

            elif is_pptx:
                prs = Presentation(template_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            for ph, val in uploads.items():
                                if ph in shape.text:
                                    shape.text = shape.text.replace(ph, str(val))
                prs.save(buffer)
                st.success("✅ PPTX generated!")
                st.download_button("📥 Download PPTX", buffer.getvalue(), file_name=final_filename + ".pptx")
